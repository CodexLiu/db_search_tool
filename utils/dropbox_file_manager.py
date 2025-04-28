import os
import tempfile
import dropbox
import dropbox.common
import io
import logging
from .dropbox_token_manager import get_access_token

# Set up logging
logger = logging.getLogger(__name__)

# Constants
MAX_DOWNLOAD_SIZE = 5 * 1024 * 1024  # 5MB absolute maximum
PARTIAL_DOWNLOAD_SIZE = 1 * 1024 * 1024  # 1MB for partial downloads

def download_file_with_size_limit(path, select_user=None, local_path=None, max_size_bytes=PARTIAL_DOWNLOAD_SIZE):
    """
    Downloads a file from Dropbox with size limit to avoid timeouts on large files.
    For files larger than MAX_DOWNLOAD_SIZE (5MB), no download is attempted.
    
    Args:
        path (str): The path of the file in Dropbox to download
        select_user (str, optional): Team member ID (e.g., "dbmid:...") to operate as.
                                     If None, operates as the admin linked to the token.
        local_path (str, optional): Local path where to save the file.
                                    If None, saves to a temporary file.
        max_size_bytes (int): Maximum file size to download (default: 1MB)
        
    Returns:
        tuple: (str: Path to the downloaded file or None if too large, 
                bool: True if file was truncated or too large, 
                int: Total file size,
                bool: True if the file was too large to download at all)
        
    Raises:
        dropbox.exceptions.ApiError: If the API request fails.
        ValueError: If token is invalid or required info cannot be determined.
        RuntimeError: For unexpected errors during the process.
    """
    access_token = get_access_token()
    if not access_token:
        raise ValueError("Failed to obtain a valid access token")

    dbx_team = dropbox.DropboxTeam(access_token)
    member_id_to_use = select_user
    is_truncated = False
    total_file_size = 0
    too_large_to_download = False

    try:
        # 1. Determine the Member ID context
        if not member_id_to_use:
            logger.info("No select_user provided, determining token admin's member ID...")
            admin_info = dbx_team.team_token_get_authenticated_admin()
            # Ensure we get the team_member_id, not just the user_id if different
            if hasattr(admin_info.admin_profile, 'team_member_id'):
                 member_id_to_use = admin_info.admin_profile.team_member_id
            else:
                 # Fallback or error, depending on API version/structure
                 raise ValueError("Could not determine team_member_id for the authenticated admin.")
        else:
            logger.info(f"Operating as specified user: {member_id_to_use}")

        # 2. Get user-specific context client
        dbx_user_context = dbx_team.as_user(member_id_to_use)

        # 3. Get account info to find the root namespace ID
        account_info = dbx_user_context.users_get_current_account()
        
        if not hasattr(account_info, 'root_info') or not hasattr(account_info.root_info, 'root_namespace_id'):
             raise ValueError("Could not determine root_namespace_id from user's account info.")
             
        root_namespace_id = account_info.root_info.root_namespace_id

        # 4. Create the final client scoped to the correct root and user
        team_path_root = dropbox.common.PathRoot.namespace_id(root_namespace_id)
        # Apply path_root first, then user context for the final client
        dbx_final_client = dbx_team.with_path_root(team_path_root).as_user(member_id_to_use)

        # 5. First, get metadata to check file size
        logger.info(f"Getting metadata for '{path}'...")
        metadata = dbx_final_client.files_get_metadata(path)
        total_file_size = metadata.size
        
        # Check if file is too large to even attempt download (> 5MB)
        if total_file_size > MAX_DOWNLOAD_SIZE:
            logger.warning(f"File size ({total_file_size} bytes) exceeds maximum limit ({MAX_DOWNLOAD_SIZE} bytes). Skipping download.")
            is_truncated = True
            too_large_to_download = True
            return None, is_truncated, total_file_size, too_large_to_download
            
        # 6. Prepare local file path
        if local_path is None:
            # Create a temporary file with the same extension as the original
            file_ext = os.path.splitext(path)[1]
            fd, temp_path = tempfile.mkstemp(suffix=file_ext)
            os.close(fd)  # Close the file descriptor
            local_path = temp_path
        else:
            # Ensure the directory exists
            os.makedirs(os.path.dirname(os.path.abspath(local_path)), exist_ok=True)
        
        # 7. Download the file (either full or partial)
        if total_file_size > max_size_bytes:
            logger.info(f"File size ({total_file_size} bytes) exceeds partial limit ({max_size_bytes} bytes). Downloading partial content...")
            is_truncated = True
            
            # Use files_download instead of files_download_to_file to get content directly
            result = dbx_final_client.files_download(path)
            # Get the file content
            _, response = result
            
            # Read only up to max_size_bytes
            content = response.content[:max_size_bytes]
            
            # Write the partial content to the file
            with open(local_path, 'wb') as f:
                f.write(content)
                
            logger.info(f"Partial download complete: {metadata.name}, downloaded: {len(content)} bytes of {total_file_size} bytes")
        else:
            logger.info(f"Downloading entire file ({total_file_size} bytes) to '{local_path}'...")
            dbx_final_client.files_download_to_file(local_path, path)
            logger.info(f"Download complete: {metadata.name}, size: {total_file_size} bytes")
        
        return local_path, is_truncated, total_file_size, too_large_to_download

    except dropbox.exceptions.AuthError as e:
        raise ValueError(f"Authentication error: {e}. Check token validity and scopes.") from e
    except dropbox.exceptions.ApiError as e:
        # Catch API errors during setup or download
        error_summary = "Unknown Error"
        if hasattr(e, 'error') and hasattr(e.error, 'get_summary'):
            error_summary = e.error.get_summary()
        elif hasattr(e, 'error'):
             error_summary = str(e.error)
        else:
             error_summary = str(e)
             
        logger.error(f"Dropbox API Error Summary: {error_summary}")
        
        # Attempt to print more details if available (e.g., for path errors)
        if hasattr(e, 'error') and hasattr(e.error, 'path') and e.error.path:
             logger.error(f"  Error details related to path lookup: {e.error.path}")
             
        # Provide context about the operation
        logger.error(f"  Error occurred during download for user {member_id_to_use} in namespace {root_namespace_id if 'root_namespace_id' in locals() else 'N/A'} for path '{path}'.")
        raise # Re-raise the original ApiError

    except Exception as e:
        # Catch any other unexpected errors
        import traceback
        logger.error(f"Unexpected Error: {e}\n{traceback.format_exc()}")
        raise RuntimeError(f"An unexpected error occurred: {e}") from e

# Update existing download_file function to use the new size-limited version
def download_file(path, select_user=None, local_path=None, max_size_bytes=PARTIAL_DOWNLOAD_SIZE):
    """
    Downloads a file from Dropbox to a temporary or specified location.
    Now uses size limiting to prevent timeouts on large files.
    
    Args:
        path (str): The path of the file in Dropbox to download
        select_user (str, optional): Team member ID (e.g., "dbmid:...") to operate as.
                                     If None, operates as the admin linked to the token.
        local_path (str, optional): Local path where to save the file.
                                    If None, saves to a temporary file.
        max_size_bytes (int): Maximum file size to download (default: 1MB)
        
    Returns:
        str: Path to the downloaded file
        
    Raises:
        dropbox.exceptions.ApiError: If the API request fails.
        ValueError: If token is invalid or required info cannot be determined.
        RuntimeError: For unexpected errors during the process.
        ValueError: If file is too large to download (>5MB)
    """
    local_path, is_truncated, total_size, too_large = download_file_with_size_limit(
        path, select_user, local_path, max_size_bytes
    )
    
    if too_large:
        size_mb = total_size / (1024 * 1024)
        raise ValueError(f"File is too large to download: {size_mb:.1f} MB exceeds the 5 MB limit")
    
    if is_truncated:
        logger.warning(f"File was truncated! Only downloaded {max_size_bytes} bytes of {total_size} bytes.")
    
    return local_path

def extract_text(file_path):
    """
    Extracts text from a file (PDF, DOCX, XLSX, PPTX).
    
    Args:
        file_path (str): Path to the file to extract text from
        
    Returns:
        str: Extracted text content
        
    Raises:
        ValueError: If the file format is not supported
        ImportError: If the required package is not installed
        RuntimeError: For unexpected errors during extraction
    """
    file_ext = os.path.splitext(file_path)[1].lower()
    
    try:
        # Extract text based on file type
        if file_ext == '.pdf':
            return _extract_text_from_pdf(file_path)
        elif file_ext == '.docx':
            return _extract_text_from_docx(file_path)
        elif file_ext == '.xlsx':
            return _extract_text_from_xlsx(file_path)
        elif file_ext == '.pptx':
            return _extract_text_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}. Supported formats: PDF, DOCX, XLSX, PPTX")
    
    except ImportError as e:
        raise ImportError(f"Required package not installed for {file_ext} extraction: {e}")
    except Exception as e:
        import traceback
        print(f"Error extracting text from {file_path}: {e}\n{traceback.format_exc()}")
        raise RuntimeError(f"Failed to extract text: {e}") from e

def _extract_text_from_pdf(file_path):
    """Extract text from a PDF file."""
    try:
        # Try using PyPDF2 first
        from PyPDF2 import PdfReader
        
        text = []
        reader = PdfReader(file_path)
        for page in reader.pages:
            text.append(page.extract_text() or "")
        
        return "\n".join(text)
    
    except ImportError:
        # Fall back to pdfminer if PyPDF2 is not available
        try:
            from pdfminer.high_level import extract_text as pdfminer_extract_text
            return pdfminer_extract_text(file_path)
        except ImportError:
            raise ImportError("Neither PyPDF2 nor pdfminer.six is installed. Install one with 'pip install PyPDF2' or 'pip install pdfminer.six'")

def _extract_text_from_docx(file_path):
    """Extract text from a DOCX file."""
    from docx import Document
    
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def _extract_text_from_xlsx(file_path):
    """Extract text from an XLSX file."""
    from openpyxl import load_workbook
    
    workbook = load_workbook(file_path, read_only=True, data_only=True)
    text = []
    
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        text.append(f"Sheet: {sheet_name}")
        
        for row in sheet.rows:
            row_text = []
            for cell in row:
                if cell.value is not None:
                    row_text.append(str(cell.value))
            if row_text:
                text.append("\t".join(row_text))
    
    return "\n".join(text)

def _extract_text_from_pptx(file_path):
    """Extract text from a PPTX file."""
    from pptx import Presentation
    
    prs = Presentation(file_path)
    text = []
    
    for i, slide in enumerate(prs.slides):
        text.append(f"Slide {i+1}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
    
    return "\n".join(text)

def download_and_extract_text(path, select_user=None, cleanup=True, max_size_bytes=PARTIAL_DOWNLOAD_SIZE):
    """
    Downloads a file from Dropbox and extracts text from it.
    
    Args:
        path (str): The path of the file in Dropbox
        select_user (str, optional): Team member ID to operate as
        cleanup (bool): Whether to delete the temporary file after extraction
        max_size_bytes (int): Maximum file size to download (default: 1MB)
        
    Returns:
        tuple: (str: Extracted text content, bool: Whether the file was truncated, int: Total file size, bool: True if file was too large)
        
    Raises:
        Various exceptions from download_file and extract_text functions
    """
    try:
        # Get metadata first to determine the file size
        local_path, is_truncated, total_size, too_large = download_file_with_size_limit(path, select_user, None, max_size_bytes)
        
        # Handle case when file is too large to download
        if too_large:
            size_mb = total_size / (1024 * 1024)
            error_message = f"This file is too large to process ({size_mb:.1f} MB exceeds the 5 MB limit). Please try a smaller file or a different approach."
            return error_message, True, total_size, True
            
        # Process the downloaded file
        text = extract_text(local_path)
        
        if is_truncated:
            text += f"\n\n[NOTE: This is a partial extraction of the file. Only the first {max_size_bytes/1024/1024:.1f} MB of {total_size/1024/1024:.1f} MB total size was processed due to file size constraints.]"
        
        if cleanup:
            try:
                os.remove(local_path)
                logger.info(f"Temporary file deleted: {local_path}")
            except Exception as e:
                logger.warning(f"Failed to delete temporary file {local_path}: {e}")
        
        return text, is_truncated, total_size, too_large
    
    except Exception as e:
        logger.error(f"Error in download_and_extract_text: {e}")
        raise 